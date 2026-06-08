#!/usr/bin/env python3
"""Generate IHG-Mamba Research Presentation PPT with speaker notes."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ── Helpers ──────────────────────────────────────────────────────────
def add_slide(prs, layout_idx=6):
    """Add blank slide."""
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])

def set_bg(slide, color_hex='FFFFFF'):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color_hex)

def add_title_box(slide, text, left, top, width, height, font_size=28, bold=True, color='1A3C6D'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor.from_string(color)
    return txBox

def add_text_box(slide, text, left, top, width, height, font_size=16, color='333333', bold=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor.from_string(color)
    p.font.bold = bold
    return txBox

def add_bullets(slide, items, left, top, width, height, font_size=15, color='333333'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor.from_string(color)
        p.space_after = Pt(8)
        p.level = 0
    return txBox

def add_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text

def add_table(slide, data, left, top, width, height):
    """data = list of lists, first row = header"""
    rows, cols = len(data), len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    for i, row in enumerate(data):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            if i == 0:
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 255)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(26, 60, 109)
            else:
                p.font.color.rgb = RGBColor(51, 51, 51)
    return table_shape

def add_divider(slide, left, top, width, color='1A3C6D'):
    shape = slide.shapes.add_shape(1, left, top, width, Pt(2))  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    shape.line.fill.background()

# ── Main ────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    W = prs.slide_width
    H = prs.slide_height
    MARGIN = Inches(0.8)
    
    # ─────────────────────────────────────────────────────────────
    # SLIDE 1: Title
    # ─────────────────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, 'F5F7FA')
    add_title_box(slide, "IHG-Mamba", Inches(1), Inches(1.8), Inches(11), Inches(1),
                  font_size=48, bold=True, color='1A3C6D')
    add_title_box(slide, "Topology-Aware Multiple Instance Learning\nfor Histopathology Image Classification",
                  Inches(1), Inches(2.9), Inches(11), Inches(1.2),
                  font_size=24, bold=False, color='4A6FA5')
    add_divider(slide, Inches(1), Inches(4.3), Inches(3))
    add_text_box(slide, "Exploring Spatial Ordering Mechanisms in Whole-Slide Image Analysis",
                 Inches(1), Inches(4.6), Inches(10), Inches(0.6), font_size=16, color='666666')
    add_text_box(slide, "Computational Pathology Lab  |  June 2026",
                 Inches(1), Inches(5.5), Inches(10), Inches(0.5), font_size=14, color='999999')
    
    add_notes(slide, """Welcome everyone. Today I'll present our work on IHG-Mamba, a topology-aware multiple instance learning framework for histopathology image classification.

Our research explores how spatial ordering of tissue patches affects the performance of MIL models in whole-slide image analysis. This is a fundamental question in computational pathology that hasn't been thoroughly investigated.""")
    
    # ─────────────────────────────────────────────────────────────
    # SLIDE 2: Background - The WSI Challenge
    # ─────────────────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, 'FFFFFF')
    add_title_box(slide, "The Whole-Slide Image Challenge", MARGIN, Inches(0.4), Inches(10), Inches(0.7),
                  font_size=28, color='1A3C6D')
    add_divider(slide, MARGIN, Inches(1.1), Inches(2.5))
    
    add_bullets(slide, [
        "• WSI resolution: 100,000 × 100,000 pixels",
        "• Each slide contains 10,000 - 100,000+ tissue patches",
        "• Cannot feed directly into CNNs/Transformers",
        "• Current solution: Multiple Instance Learning (MIL)",
    ], Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5), font_size=16)
    
    add_bullets(slide, [
        "• MIL treats patches as an unordered 'bag'",
        "• Aggregation via attention or pooling",
        "• Problem: spatial structure is ignored!",
        "• Tissue morphology has spatial continuity",
    ], Inches(7), Inches(1.5), Inches(5.5), Inches(2.5), font_size=16)
    
    add_text_box(slide, "Key Question: Does patch ordering matter for MIL?",
                 Inches(2), Inches(5), Inches(9), Inches(0.6), font_size=18, bold=True, color='C0392B')
    
    add_notes(slide, """Let me start with the background. Whole-slide images are extremely high resolution - typically 100,000 by 100,000 pixels. Each slide contains tens of thousands of tissue patches.

Current approaches use Multiple Instance Learning, or MIL, which treats patches as an unordered bag. But here's the key problem: tissue morphology has spatial continuity. A tumor region transitions gradually into normal tissue. By ignoring spatial structure, we're throwing away potentially valuable information.

This leads to our key question: does patch ordering actually matter for MIL performance?""")
    
    # ─────────────────────────────────────────────────────────────
    # SLIDE 3: Research Question & Hypothesis
    # ─────────────────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, 'F5F7FA')
    add_title_box(slide, "Research Question & Hypothesis", MARGIN, Inches(0.4), Inches(10), Inches(0.7),
                  font_size=28, color='1A3C6D')
    add_divider(slide, MARGIN, Inches(1.1), Inches(2.5))
    
    add_text_box(slide, "Research Question:", Inches(0.8), Inches(1.5), Inches(10), Inches(0.5),
                 font_size=18, bold=True, color='1A3C6D')
    add_text_box(slide, "How does the spatial ordering of tissue patches affect\nMIL-based histopathology classification performance?",
                 Inches(1.2), Inches(2.1), Inches(10), Inches(0.8), font_size=16, color='333333')
    
    add_text_box(slide, "Hypothesis:", Inches(0.8), Inches(3.2), Inches(10), Inches(0.5),
                 font_size=18, bold=True, color='1A3C6D')
    add_text_box(slide, "Hilbert space-filling curve ordering can preserve spatial\ncontinuity and improve classification accuracy.",
                 Inches(1.2), Inches(3.8), Inches(10), Inches(0.8), font_size=16, color='333333')
    
    add_text_box(slide, "Approach:", Inches(0.8), Inches(4.9), Inches(10), Inches(0.5),
                 font_size=18, bold=True, color='1A3C6D')
    add_bullets(slide, [
        "• Compare: Hilbert order vs Raw/Raster order vs Random order",
        "• Use state-space models (Mamba) for sequence modeling",
        "• Evaluate on LUAD/LUSC classification (1052 WSIs)",
    ], Inches(1.2), Inches(5.5), Inches(10), Inches(1.5), font_size=15)
    
    add_notes(slide, """Our research question is straightforward: how does spatial ordering affect MIL performance?

Our hypothesis is that Hilbert space-filling curve ordering can preserve spatial continuity. The Hilbert curve is a fractal curve that visits every point in a 2D grid while maintaining spatial locality - nearby points in 1D are also nearby in 2D.

Our approach is to compare three ordering strategies: Hilbert, raw raster scan, and random permutation. We use Mamba, a state-space model, for sequence modeling because it can efficiently handle long sequences. And we evaluate on LUAD versus LUSC classification with 1052 whole-slide images.""")
    
    # ─────────────────────────────────────────────────────────────
    # SLIDE 4: IHG-Mamba Architecture
    # ─────────────────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, 'FFFFFF')
    add_title_box(slide, "IHG-Mamba Architecture", MARGIN, Inches(0.4), Inches(10), Inches(0.7),
                  font_size=28, color='1A3C6D')
    add_divider(slide, MARGIN, Inches(1.1), Inches(2.5))
    
    # Architecture pipeline
    steps = [
        ("1. Hilbert\nOrdering", "Preserve spatial\nlocality"),
        ("2. Local\nMamba", "Model local\ntissue context"),
        ("3. Super-node\nAvgPool", "Aggregate into\ntissue regions"),
        ("4. Global\nMamba", "Capture long-range\ndependencies"),
        ("5. Attention\nReadout", "Weighted slide-level\nprediction"),
    ]
    
    for i, (title, desc) in enumerate(steps):
        x = Inches(0.5 + i * 2.5)
        y = Inches(1.8)
        # Box
        shape = slide.shapes.add_shape(1, x, y, Inches(2.2), Inches(1.8))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(26, 60, 109) if i % 2 == 0 else RGBColor(74, 111, 165)
        shape.line.fill.background()
        
        txBox = slide.shapes.add_textbox(x, y + Inches(0.2), Inches(2.2), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        
        txBox2 = slide.shapes.add_textbox(x, y + Inches(1.0), Inches(2.2), Inches(0.6))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(200, 220, 255)
        p2.alignment = PP_ALIGN.CENTER
        
        # Arrow
        if i < len(steps) - 1:
            arrow_x = x + Inches(2.2)
            txArrow = slide.shapes.add_textbox(arrow_x, y + Inches(0.7), Inches(0.3), Inches(0.4))
            tf_a = txArrow.text_frame
            p_a = tf_a.paragraphs[0]
            p_a.text = "→"
            p_a.font.size = Pt(20)
            p_a.font.bold = True
            p_a.font.color.rgb = RGBColor(26, 60, 109)
    
    add_text_box(slide, "Key insight: Hilbert ordering enables sequence models to capture spatial dependencies",
                 Inches(1), Inches(5.5), Inches(11), Inches(0.6), font_size=16, bold=True, color='C0392B')
    
    add_text_box(slide, "Input: 10,000+ patches (1024-dim UNI features) → Output: slide-level prediction (LUAD/LUSC)",
                 Inches(1), Inches(6.2), Inches(11), Inches(0.5), font_size=13, color='666666')
    
    add_notes(slide, """Now let me explain our IHG-Mamba architecture. It's a five-stage pipeline.

First, we order all tissue patches using the Hilbert space-filling curve. This preserves spatial locality - patches that are close in 2D space remain close in the 1D sequence.

Second, we apply Local Mamba to model dependencies within local tissue neighborhoods. Mamba is a state-space model that's much more efficient than Transformers for long sequences.

Third, we aggregate local features into super-nodes using average pooling. This reduces the sequence length while preserving regional information.

Fourth, we apply Global Mamba to capture long-range dependencies between super-nodes.

Finally, attention readout produces the slide-level prediction.

The key insight is that by preserving spatial ordering, the sequence model can actually learn meaningful spatial patterns in tissue morphology.""")
    
    # ─────────────────────────────────────────────────────────────
    # SLIDE 5: Experimental Design
    # ─────────────────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, 'F5F7FA')
    add_title_box(slide, "Experimental Design", MARGIN, Inches(0.4), Inches(10), Inches(0.7),
                  font_size=28, color='1A3C6D')
    add_divider(slide, MARGIN, Inches(1.1), Inches(2.5))
    
    # Left: Dataset
    add_text_box(slide, "Dataset", Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
                 font_size=18, bold=True, color='1A3C6D')
    dataset_data = [
        ['Task', 'Samples', 'Classes', 'Feature'],
        ['LUAD/LUSC', '1052 WSI', '2 (binary)', 'UNI 1024-dim'],
        ['OV Tumor/Normal', '390 WSI', '2 (binary)', 'UNI 1024-dim'],
    ]
    add_table(slide, dataset_data, Inches(0.8), Inches(2.1), Inches(5.5), Inches(1.2))
    
    # Right: Methods
    add_text_box(slide, "Methods Compared", Inches(7), Inches(1.5), Inches(5), Inches(0.5),
                 font_size=18, bold=True, color='1A3C6D')
    method_data = [
        ['Method', 'Type'],
        ['IHG-Mamba (ours)', 'Topology-aware MIL'],
        ['CLAM', 'Attention MIL'],
        ['ABMIL', 'Attention MIL'],
    ]
    add_table(slide, method_data, Inches(7), Inches(2.1), Inches(5.5), Inches(1.2))
    
    # Bottom: Ablations
    add_text_box(slide, "Ablation Studies", Inches(0.8), Inches(4), Inches(10), Inches(0.5),
                 font_size=18, bold=True, color='1A3C6D')
    add_bullets(slide, [
        "• Ordering: Hilbert vs Raw/Raster vs RandomPerm",
        "• Pooling: AvgPool vs BRPool vs BP-Pool",
        "• Readout: Simple Attention vs Gated Attention",
        "• Evaluation: 5-fold cross-validation, paired statistical tests",
    ], Inches(1.2), Inches(4.6), Inches(10), Inches(2), font_size=15)
    
    add_notes(slide, """Our experimental setup uses two datasets. The main one is LUAD versus LUSC classification with 1052 whole-slide images from TCGA. We also validate on an OV tumor versus normal task with 390 slides.

We compare against CLAM and ABMIL, two standard MIL baselines.

Our ablation studies systematically vary three components: the ordering strategy, the pooling mechanism, and the attention readout. All experiments use 5-fold cross-validation with paired statistical tests.""")
    
    # ─────────────────────────────────────────────────────────────
    # SLIDE 6: Main Results
    # ─────────────────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, 'FFFFFF')
    add_title_box(slide, "Main Results: LUAD/LUSC Classification", MARGIN, Inches(0.4), Inches(10), Inches(0.7),
                  font_size=28, color='1A3C6D')
    add_divider(slide, MARGIN, Inches(1.1), Inches(2.5))
    
    results_data = [
        ['Method', 'Test AUC', 'Test Accuracy', 'Note'],
        ['IHG-Mamba (Hilbert+Avg)', '0.9808±0.010', '0.9448±0.023', 'Best'],
        ['CLAM-SB', '0.9786±0.013', '0.9333±0.030', 'Baseline'],
        ['ABMIL (raw order)', '0.9725±0.010', '0.9240±0.012', 'Baseline'],
        ['IHG-Mamba (Raw)', '0.9752±0.014', '0.9257±0.040', 'No Hilbert'],
        ['IHG-Mamba (RandomPerm)', '0.9608±0.018', '0.9143±0.034', 'Random order'],
    ]
    add_table(slide, results_data, Inches(0.8), Inches(1.5), Inches(11.5), Inches(2.5))
    
    add_text_box(slide, "Key findings:", Inches(0.8), Inches(4.3), Inches(10), Inches(0.5),
                 font_size=18, bold=True, color='1A3C6D')
    add_bullets(slide, [
        "✓ IHG-Mamba achieves best AUC (0.9808) among all methods",
        "✓ Significant improvement over RandomPerm (p=0.025*)",
        "✓ Competitive with CLAM and ABMIL baselines",
        "✓ Multi-seed stable: seed1=0.9808, seed2=0.9793",
    ], Inches(1.2), Inches(4.9), Inches(10), Inches(2), font_size=15)
    
    add_notes(slide, """Here are our main results on LUAD versus LUSC classification.

Our IHG-Mamba with Hilbert ordering and average pooling achieves the best AUC of 0.9808, which is competitive with CLAM at 0.9786 and better than ABMIL at 0.9725.

Most importantly, the difference between Hilbert and RandomPerm ordering is statistically significant with p equals 0.025. This directly answers our research question: yes, spatial ordering matters.

The results are also stable across random seeds - we get 0.9808 with seed 1 and 0.9793 with seed 2.""")
    
    # ─────────────────────────────────────────────────────────────
    # SLIDE 7: Ordering Ablation
    # ─────────────────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, 'F5F7FA')
    add_title_box(slide, "Ordering Matters: Hilbert vs Random", MARGIN, Inches(0.4), Inches(10), Inches(0.7),
                  font_size=28, color='1A3C6D')
    add_divider(slide, MARGIN, Inches(1.1), Inches(2.5))
    
    # Left: AUC comparison
    add_text_box(slide, "Classification Performance", Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
                 font_size=18, bold=True, color='1A3C6D')
    
    order_data = [
        ['Ordering', 'Test AUC', 'Δ vs Hilbert'],
        ['Hilbert', '0.9808±0.010', '—'],
        ['Raw/Raster', '0.9752±0.014', '-0.0056 (ns)'],
        ['RandomPerm', '0.9608±0.018', '-0.0200 (p=0.025*)'],
    ]
    add_table(slide, order_data, Inches(0.8), Inches(2.1), Inches(5.5), Inches(1.5))
    
    # Right: Spatial metrics
    add_text_box(slide, "Spatial Continuity Metrics", Inches(7), Inches(1.5), Inches(5), Inches(0.5),
                 font_size=18, bold=True, color='1A3C6D')
    
    spatial_data = [
        ['Ordering', 'Mean Jump', 'Tear Rate'],
        ['Hilbert', '279', '0.18%'],
        ['Raw/Raster', '540', '1.29%'],
        ['RandomPerm', '18,061', '99.17%'],
    ]
    add_table(slide, spatial_data, Inches(7), Inches(2.1), Inches(5.5), Inches(1.5))
    
    add_text_box(slide, "Conclusion: Spatial continuity correlates with classification performance",
                 Inches(1), Inches(5.5), Inches(11), Inches(0.6), font_size=17, bold=True, color='C0392B')
    
    add_notes(slide, """This slide shows the core finding of our work.

On the left, we see classification performance by ordering strategy. Hilbert ordering gives the best AUC at 0.9808. Random permutation drops significantly to 0.9608, with p equals 0.025.

On the right, we see spatial continuity metrics. Mean jump distance measures how far we move between consecutive patches in the sequence. Hilbert has the smallest jump at 279 pixels, while RandomPerm jumps 18,000 pixels on average.

Tear rate measures how often we cross tissue boundaries. Hilbert has only 0.18% tear rate versus 99% for random ordering.

The conclusion is clear: spatial continuity directly correlates with classification performance. This validates our hypothesis that preserving spatial structure matters for MIL.""")
    
    # ─────────────────────────────────────────────────────────────
    # SLIDE 8: Additional Results
    # ─────────────────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, 'FFFFFF')
    add_title_box(slide, "Additional Validation", MARGIN, Inches(0.4), Inches(10), Inches(0.7),
                  font_size=28, color='1A3C6D')
    add_divider(slide, MARGIN, Inches(1.1), Inches(2.5))
    
    # OV results
    add_text_box(slide, "OV Tumor vs Normal Classification", Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
                 font_size=18, bold=True, color='1A3C6D')
    add_bullets(slide, [
        "• Task: Tumor vs Normal tissue (390 WSIs)",
        "• Result: AUC = 1.0, Accuracy = 99.7%",
        "• 5-fold CV, all folds AUC = 1.0",
        "• Note: Task is relatively simple (morphological difference)",
    ], Inches(1.2), Inches(2.1), Inches(5), Inches(2), font_size=14)
    
    # Efficiency
    add_text_box(slide, "Efficiency Analysis", Inches(7), Inches(1.5), Inches(5), Inches(0.5),
                 font_size=18, bold=True, color='1A3C6D')
    eff_data = [
        ['Tokens', 'AUC', 'Relative'],
        ['2500', '0.9749', '99.4%'],
        ['5000', '0.9776', '99.7%'],
        ['Full', '0.9808', '100%'],
    ]
    add_table(slide, eff_data, Inches(7), Inches(2.1), Inches(5), Inches(1.5))
    
    # Survival
    add_text_box(slide, "Survival Analysis (TCGA-LUAD)", Inches(0.8), Inches(4.5), Inches(5), Inches(0.5),
                 font_size=18, bold=True, color='1A3C6D')
    add_bullets(slide, [
        "• Task: Overall survival prediction (454 WSIs)",
        "• Best C-index: 0.6117±0.052",
        "• Task ceiling ~0.61 (limited by sample size)",
        "• IHG ordering shows no significant improvement",
    ], Inches(1.2), Inches(5.1), Inches(5), Inches(1.5), font_size=14)
    
    add_notes(slide, """Let me briefly show some additional results.

For the OV tumor versus normal task, we achieve perfect AUC of 1.0. However, this task is relatively simple because tumor and normal tissue have very different morphology.

For efficiency, we find that using only 2500 tokens achieves 99.4% of full performance. This is important for practical deployment.

For survival analysis, the C-index is around 0.61, which is typical for this task. However, we don't see improvement from Hilbert ordering here, likely because survival prediction depends more on molecular features than spatial structure.""")
    
    # ─────────────────────────────────────────────────────────────
    # SLIDE 9: Key Contributions
    # ─────────────────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, 'F5F7FA')
    add_title_box(slide, "Key Contributions", MARGIN, Inches(0.4), Inches(10), Inches(0.7),
                  font_size=28, color='1A3C6D')
    add_divider(slide, MARGIN, Inches(1.1), Inches(2.5))
    
    contributions = [
        ("1", "First systematic study of spatial ordering in MIL",
         "Demonstrated that patch ordering significantly affects classification (p=0.025)"),
        ("2", "Topology-aware architecture IHG-Mamba",
         "Hilbert curve + Local/Global Mamba for spatial sequence modeling"),
        ("3", "Spatial continuity metrics",
         "Introduced jump distance and tear rate to quantify ordering quality"),
        ("4", "Practical efficiency",
         "99.4% performance with only 2500 tokens (25% of full)"),
    ]
    
    for i, (num, title, desc) in enumerate(contributions):
        y = Inches(1.5 + i * 1.4)
        # Number circle
        shape = slide.shapes.add_shape(9, Inches(0.8), y, Inches(0.5), Inches(0.5))  # oval
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(26, 60, 109)
        shape.line.fill.background()
        
        txNum = slide.shapes.add_textbox(Inches(0.8), y + Inches(0.05), Inches(0.5), Inches(0.4))
        tf_n = txNum.text_frame
        p_n = tf_n.paragraphs[0]
        p_n.text = num
        p_n.font.size = Pt(16)
        p_n.font.bold = True
        p_n.font.color.rgb = RGBColor(255, 255, 255)
        p_n.alignment = PP_ALIGN.CENTER
        
        add_text_box(slide, title, Inches(1.5), y, Inches(10), Inches(0.4),
                     font_size=18, bold=True, color='1A3C6D')
        add_text_box(slide, desc, Inches(1.5), y + Inches(0.45), Inches(10), Inches(0.5),
                     font_size=14, color='555555')
    
    add_notes(slide, """Let me summarize our key contributions.

First, this is the first systematic study showing that spatial ordering matters for MIL in histopathology. The p-value of 0.025 provides statistical evidence.

Second, we propose IHG-Mamba, a topology-aware architecture that combines Hilbert ordering with state-space models for efficient spatial modeling.

Third, we introduce spatial continuity metrics - jump distance and tear rate - that can quantify the quality of any ordering strategy.

Fourth, we show practical efficiency gains: 99.4% performance with only a quarter of the tokens.""")
    
    # ─────────────────────────────────────────────────────────────
    # SLIDE 10: Limitations & Future Work
    # ─────────────────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, 'FFFFFF')
    add_title_box(slide, "Limitations & Future Directions", MARGIN, Inches(0.4), Inches(10), Inches(0.7),
                  font_size=28, color='1A3C6D')
    add_divider(slide, MARGIN, Inches(1.1), Inches(2.5))
    
    add_text_box(slide, "Current Limitations", Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
                 font_size=18, bold=True, color='C0392B')
    add_bullets(slide, [
        "• Survival analysis: no significant improvement",
        "• Task ceiling ~0.61 C-index (limited data)",
        "• BRPool/BP-Pool: no additional benefit",
        "• Single dataset validation (TCGA)",
    ], Inches(1.2), Inches(2.1), Inches(5), Inches(2), font_size=14)
    
    add_text_box(slide, "Future Directions", Inches(7), Inches(1.5), Inches(5), Inches(0.5),
                 font_size=18, bold=True, color='27AE60')
    add_bullets(slide, [
        "• External validation on independent cohorts",
        "• Multi-class cancer subtyping",
        "• Integration with clinical data",
        "• Pathologist review of attention heatmaps",
        "• Extension to other imaging modalities",
    ], Inches(7.4), Inches(2.1), Inches(5), Inches(2.5), font_size=14)
    
    add_text_box(slide, "Open Questions", Inches(0.8), Inches(4.5), Inches(11), Inches(0.5),
                 font_size=18, bold=True, color='1A3C6D')
    add_bullets(slide, [
        "• Why does ordering not help survival prediction?",
        "• Is there an optimal ordering for specific cancer types?",
        "• Can we learn the ordering end-to-end?",
    ], Inches(1.2), Inches(5.1), Inches(10), Inches(1.5), font_size=14)
    
    add_notes(slide, """Now let me discuss limitations and future work.

Current limitations include that we don't see improvement for survival analysis, which has a task ceiling around 0.61 C-index. We also haven't validated on external datasets yet.

Future directions include external validation on independent cohorts, multi-class cancer subtyping, and integration with clinical data. We also plan to have pathologists review our attention heatmaps to verify biological relevance.

Some open questions remain: why doesn't ordering help survival prediction? Is there an optimal ordering that depends on cancer type? And can we learn the ordering end-to-end instead of using a fixed Hilbert curve?""")
    
    # ─────────────────────────────────────────────────────────────
    # SLIDE 11: Summary & Thank You
    # ─────────────────────────────────────────────────────────────
    slide = add_slide(prs)
    set_bg(slide, 'F5F7FA')
    add_title_box(slide, "Summary", MARGIN, Inches(0.4), Inches(10), Inches(0.7),
                  font_size=28, color='1A3C6D')
    add_divider(slide, MARGIN, Inches(1.1), Inches(2.5))
    
    add_text_box(slide, "Main Findings:", Inches(0.8), Inches(1.5), Inches(10), Inches(0.5),
                 font_size=20, bold=True, color='1A3C6D')
    
    findings = [
        "1. Spatial ordering significantly affects MIL performance (p=0.025*)",
        "2. Hilbert curve preserves tissue continuity (tear rate: 0.18% vs 99%)",
        "3. IHG-Mamba achieves competitive AUC (0.9808) on LUAD/LUSC",
        "4. Practical efficiency: 99.4% performance with 25% tokens",
    ]
    add_bullets(slide, findings, Inches(1.2), Inches(2.2), Inches(10), Inches(2.5), font_size=17)
    
    add_text_box(slide, "Significance:", Inches(0.8), Inches(4.5), Inches(10), Inches(0.5),
                 font_size=20, bold=True, color='1A3C6D')
    add_text_box(slide, "This work reveals a fundamental mechanism in pathology AI:\nspatial structure matters for whole-slide image understanding.",
                 Inches(1.2), Inches(5.1), Inches(10), Inches(0.8), font_size=16, color='333333')
    
    add_text_box(slide, "Thank you!  Questions?", Inches(3), Inches(6.2), Inches(7), Inches(0.6),
                 font_size=24, bold=True, color='1A3C6D')
    
    add_notes(slide, """To summarize, our work demonstrates four key findings.

First, spatial ordering significantly affects MIL performance with statistical significance.

Second, Hilbert curve ordering preserves tissue continuity with only 0.18% tear rate compared to 99% for random ordering.

Third, our IHG-Mamba achieves competitive AUC of 0.9808 on LUAD versus LUSC classification.

Fourth, we achieve practical efficiency gains with 99.4% performance using only 25% of tokens.

The significance of this work is that it reveals a fundamental mechanism in pathology AI: spatial structure matters for whole-slide image understanding.

Thank you for your attention. I'm happy to take questions.""")
    
    # ── Save ─────────────────────────────────────────────────────
    output_dir = os.path.join(BASE, 'reports')
    os.makedirs(output_dir, exist_ok=True)
    output_path = os.path.join(output_dir, 'IHG_Mamba_Presentation.pptx')
    prs.save(output_path)
    print(f"PPTX saved: {output_path}")
    print(f"Slides: {len(prs.slides)}")


BASE = '/home/a255372639/projects/IHG-Mamba'
if __name__ == '__main__':
    main()
